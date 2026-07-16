"""PPTX 안의 모든 삽입 이미지를 슬라이드 번호 기준으로 자동 추출합니다.
사용법: python extract_ppt_images.py Apple_운영_라인업.pptx assets/ppt-images
"""
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import sys

if len(sys.argv) < 2:
    raise SystemExit("사용법: python extract_ppt_images.py <pptx 경로> [출력 폴더]")
pptx_path=Path(sys.argv[1])
out=Path(sys.argv[2] if len(sys.argv)>2 else "assets/ppt-images")
out.mkdir(parents=True,exist_ok=True)
prs=Presentation(pptx_path)
count=0
for slide_no,slide in enumerate(prs.slides,1):
    picture_no=0
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        picture_no+=1
        ext=shape.image.ext or "png"
        filename=out/f"slide-{slide_no:02d}-image-{picture_no:02d}.{ext}"
        filename.write_bytes(shape.image.blob)
        count+=1
print(f"{count}개 이미지를 {out}에 추출했습니다.")
