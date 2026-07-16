from pptx import Presentation
from pathlib import Path
from collections import defaultdict
from PIL import Image
from io import BytesIO
import json, re, shutil, os

PPT = Path('/mnt/data/2607)Apple_운영_라인업 (1).pptx')
ROOT = Path('/mnt/data/apple-consult-full')
ASSET = ROOT/'assets'/'products'
ASSET.mkdir(parents=True, exist_ok=True)
prs = Presentation(str(PPT))

def clean(s):
    if s is None: return ''
    s = str(s).replace('\x0b',' ').replace('\n',' / ').replace('”','"').replace('“','"')
    s = re.sub(r'\s+',' ',s).strip()
    s = re.sub(r'\s*/\s*',' / ',s)
    return s

def money(s):
    s=clean(s).replace(' ','').replace(',','')
    m=re.search(r'\d+',s)
    return int(m.group()) if m else None

def slugify(s):
    s=s.lower().replace('+','plus')
    s=re.sub(r'[^a-z0-9가-힣]+','-',s).strip('-')
    return s

def rows(slide_no, table_idx=0):
    tables=[sh.table for sh in prs.slides[slide_no-1].shapes if sh.has_table]
    t=tables[table_idx]
    return [[clean(t.cell(r,c).text) for c in range(len(t.columns))] for r in range(len(t.rows))]

def ffill_table(data, cols):
    last={c:'' for c in cols}
    out=[]
    for row in data:
        row=row[:]
        for c in cols:
            if c < len(row):
                if row[c]: last[c]=row[c]
                else: row[c]=last[c]
        out.append(row)
    return out

def normalize_family(s):
    s=clean(s).replace(' / ',' ')
    s=re.sub(r'ProMax','Pro Max',s,flags=re.I)
    s=re.sub(r'17Pro','17 Pro',s,flags=re.I)
    s=re.sub(r'16Pro','16 Pro',s,flags=re.I)
    s=re.sub(r'\s+',' ',s).strip()
    return s

def add_product(store, category, family, variant, image=None, description='', source_url='', source_label='Apple 공식 비교'):
    family=normalize_family(family)
    key=f'{category}::{family}'
    if key not in store:
        store[key]={
            'id': slugify(key), 'category': category, 'family': family,
            'image': image or '', 'description': description,
            'sourceUrl': source_url, 'sourceLabel': source_label,
            'variants': []
        }
    if image and not store[key]['image']: store[key]['image']=image
    clean_variant={}
    for k,v in variant.items():
        if v in ('',None,'-'): continue
        if k=='modelCode' and isinstance(v,str):
            v=re.sub(r'\s*/\s*','/',v).replace(' ','').upper()
        clean_variant[k]=v
    store[key]['variants'].append(clean_variant)

# image extraction mapping: product family -> (slide, picture index)
image_map = {
 'iPhone 16e':(3,1), 'iPhone 17e':(3,2), 'iPhone 17':(4,1), 'iPhone Air':(4,2),
 'iPhone 17 Pro':(5,1), 'iPhone 17 Pro Max':(5,2), 'iPhone 16':(7,1), 'iPhone 16 Plus':(7,2),
 'iPhone 16 Pro':(8,1), 'iPhone 16 Pro Max':(8,2),
 'iPad Pro 11':(10,1), 'iPad Pro 13':(10,2), 'iPad Air 11':(11,1), 'iPad Air 13':(11,2),
 'iPad A16':(12,1), 'iPad mini':(13,1),
 'Apple Watch Series 11':(15,1), 'Apple Watch SE3':(17,1), 'Apple Watch Ultra3':(18,1),
 'AirPods Pro 3':(20,1), 'AirPods Pro 2':(20,2), 'AirPods 4(ANC)':(20,4), 'AirPods 4':(20,4), 'AirPods Max':(20,3),
 'iMac':(22,1), 'MacBook Air':(26,1), 'MacBook Pro':(26,3), 'Mac mini':(27,3), 'Mac Studio':(28,1),
 'Apple Pencil Pro':(33,1), 'Apple Pencil 1세대':(33,2), 'Apple Pencil USB-C':(33,3),
 '매직키보드':(33,4), '스마트폴리오':(34,1), '워치밴드':(35,1),
 'Mac 액세서리':(38,1), '충전기':(39,2), '케이블':(40,1), 'AppleCare+':(6,1)
}

def get_picture(slide_no, pic_idx):
    pics=[sh for sh in prs.slides[slide_no-1].shapes if sh.shape_type==13]
    if pic_idx-1 >= len(pics): return None
    return pics[pic_idx-1]

image_files={}
for family,(sn,pi) in image_map.items():
    sh=get_picture(sn,pi)
    if not sh: continue
    blob=sh.image.blob
    try:
        im=Image.open(BytesIO(blob)).convert('RGBA')
        # crop transparent/white margin lightly
        bg=Image.new('RGBA',im.size,(255,255,255,0))
        # preserve original, webp smaller
        fn=f'{slugify(family)}.webp'
        im.save(ASSET/fn,'WEBP',quality=90,method=6)
        image_files[family]=f'assets/products/{fn}'
    except Exception:
        ext=sh.image.ext or 'png'; fn=f'{slugify(family)}.{ext}'
        (ASSET/fn).write_bytes(blob); image_files[family]=f'assets/products/{fn}'

store={}
source_urls={
 'iPhone':'https://www.apple.com/kr/iphone/compare/',
 'iPad':'https://www.apple.com/kr/ipad/compare/',
 'Apple Watch':'https://www.apple.com/kr/watch/compare/',
 'AirPods':'https://www.apple.com/kr/airpods/compare/',
 'Mac':'https://www.apple.com/kr/mac/compare/',
 'Accessory':'https://www.apple.com/kr/shop/accessories/all',
 'AppleCare+':'https://www.apple.com/kr/support/products/'
}

descriptions={
 'iPhone 17e':'실속형 iPhone. 가격과 기본 사용성을 중시하는 고객에게 적합합니다.',
 'iPhone 17':'표준형 iPhone. 화면, 성능, 카메라의 균형을 원하는 고객에게 적합합니다.',
 'iPhone Air':'큰 화면과 얇고 가벼운 사용감을 우선하는 고객에게 적합합니다.',
 'iPhone 17 Pro':'고성능과 프로급 촬영 기능을 원하는 고객을 위한 모델입니다.',
 'iPhone 17 Pro Max':'가장 큰 화면과 최대 용량을 원하는 고객에게 적합합니다.',
 'iPad A16':'학습, 영상, 기본 문서 작업 중심의 입문형 iPad입니다.',
 'iPad Air 11':'휴대성과 성능의 균형이 좋은 범용 iPad입니다.',
 'iPad Air 13':'큰 화면에서 필기, 영상, 문서 작업을 원하는 고객에게 적합합니다.',
 'iPad Pro 11':'휴대 가능한 프로 작업용 iPad입니다.',
 'iPad Pro 13':'대화면과 고성능을 원하는 전문 작업 고객에게 적합합니다.',
 'iPad mini':'작고 가벼운 휴대성을 최우선으로 하는 고객에게 적합합니다.',
 'Apple Watch SE3':'기본 알림, 운동, 건강 관리와 가격을 중시하는 고객에게 적합합니다.',
 'Apple Watch Series 11':'상시표시 화면과 강화된 건강 기능을 원하는 고객에게 적합합니다.',
 'Apple Watch Ultra3':'아웃도어와 전문 스포츠 활용을 원하는 고객에게 적합합니다.',
 'AirPods 4':'편안한 오픈형 착용감과 기본 기능을 원하는 고객에게 적합합니다.',
 'AirPods 4(ANC)':'오픈형 착용감은 유지하면서 노이즈 캔슬링이 필요한 고객에게 적합합니다.',
 'AirPods Pro 3':'차음성, 강한 노이즈 캔슬링, 운동 활용을 중시하는 고객에게 적합합니다.',
 'AirPods Max':'오버이어 착용감과 몰입형 청취를 원하는 고객에게 적합합니다.',
 'MacBook Neo':'가벼운 일상 작업과 합리적인 가격을 우선하는 고객에게 적합합니다.',
 'MacBook Air':'휴대성, 긴 사용 시간, 일상 및 업무 성능의 균형이 좋은 모델입니다.',
 'MacBook Pro':'고성능 작업, 영상 편집, 개발 등 전문 작업을 위한 모델입니다.',
 'iMac':'본체와 디스플레이가 결합된 깔끔한 데스크톱 환경에 적합합니다.',
 'Mac mini':'기존 모니터와 주변기기를 활용하는 소형 데스크톱입니다.',
 'Mac Studio':'고성능 데스크톱 작업 환경이 필요한 전문가용 모델입니다.'
}

# iPhone
for sn in [3,4,5,7,8]:
    data=ffill_table(rows(sn)[1:],[0,3,4,5,6])
    for r in data:
        fam=normalize_family(r[0]);
        if not fam or not r[8]: continue
        add_product(store,'iPhone',fam,{
            'display':r[3], 'ram':r[4], 'storage':r[5], 'price':money(r[6]), 'color':r[7], 'modelCode':r[8]
        },image_files.get(fam),descriptions.get(fam,''),source_urls['iPhone'])

# iPad
for sn in [10,11,12,13]:
    data=ffill_table(rows(sn)[2:],[0,3,4,5,9,10])
    for r in data:
        fam=normalize_family(r[0]);
        if not fam: continue
        for conn,pc,cc in [('Wi-Fi',4,7),('Cellular',5,8)]:
            code=r[cc] if cc < len(r) else ''
            if code and code!='-':
                add_product(store,'iPad',fam,{
                    'storage':r[3], 'connectivity':conn, 'price':money(r[pc]), 'color':r[6], 'modelCode':code,
                    'pencil':r[9], 'keyboard':r[10]
                },image_files.get(fam),descriptions.get(fam,''),source_urls['iPad'])

# Watch Series/SE
for sn in [15,16,17]:
    data=ffill_table(rows(sn)[2:],[0,2,3,5,6])
    for r in data:
        raw=normalize_family(r[0])
        fam={'Series 11':'Apple Watch Series 11','SE3':'Apple Watch SE3'}.get(raw,raw)
        if not fam: continue
        base={'size':r[2],'material':r[3],'band':r[4],'caseColor':r[7],'bandColor':r[8]}
        if r[9] and r[9]!='-': add_product(store,'Apple Watch',fam,{**base,'connectivity':'GPS','price':money(r[5]),'modelCode':r[9]},image_files.get(fam),descriptions.get(fam,''),source_urls['Apple Watch'])
        if r[10] and r[10]!='-': add_product(store,'Apple Watch',fam,{**base,'connectivity':'GPS + Cellular','price':money(r[6]),'modelCode':r[10]},image_files.get(fam),descriptions.get(fam,''),source_urls['Apple Watch'])
# Ultra
for r in ffill_table(rows(18)[2:],[0,2,3,5]):
    fam='Apple Watch Ultra3';
    if r[8] and r[8]!='-':
        add_product(store,'Apple Watch',fam,{'size':r[2],'material':r[3],'band':r[4],'connectivity':'GPS + Cellular','price':money(r[5]),'caseColor':r[6],'bandColor':r[7],'modelCode':r[8]},image_files.get(fam),descriptions.get(fam,''),source_urls['Apple Watch'])

# AirPods
for r in ffill_table(rows(20)[1:],[0,4]):
    fam=normalize_family(r[0]);
    if not fam or not r[3]: continue
    add_product(store,'AirPods',fam,{'color':r[2],'modelCode':r[3],'price':money(r[4]),'charging':r[5],'anc':'지원' if r[6]=='O' else '미지원','port':r[7]},image_files.get(fam),descriptions.get(fam,''),source_urls['AirPods'])

# Mac generic parser
# iMac
for r in ffill_table(rows(22)[1:],[0,2,4,5,6,7,10]):
    fam='iMac'
    add_product(store,'Mac',fam,{'display':r[2],'cpu':r[4],'gpu':r[5],'ram':r[6],'storage':r[7],'color':r[8],'modelCode':r[9],'price':money(r[10])},image_files.get(fam),descriptions.get(fam,''),source_urls['Mac'])
# Neo
for r in ffill_table(rows(23)[1:],[0,3,4,5,8]):
    fam='MacBook Neo'
    add_product(store,'Mac',fam,{'display':r[3],'ram':r[4],'storage':r[5],'color':re.sub(r'\s+','',r[6]),'modelCode':r[7],'price':money(r[8]),'chip':'A18 Pro'},'',descriptions.get(fam,''),source_urls['Mac'])
# Air
for r in ffill_table(rows(24)[1:],[0,3,4,5,6,7,10]):
    fam='MacBook Air'
    add_product(store,'Mac',fam,{'display':r[3],'cpu':r[4],'gpu':r[5],'ram':r[6],'storage':r[7],'color':r[8],'modelCode':r[9],'price':money(r[10]),'chip':'M5'},image_files.get(fam),descriptions.get(fam,''),source_urls['Mac'])
# Pro slides 25,26
for sn in [25,26]:
    rr=rows(sn)[1:]
    # slide25 order SSD/RAM, slide26 RAM/SSD
    if sn==25:
        data=ffill_table(rr,[0,2,4,5,6,7,10])
        for r in data:
            add_product(store,'Mac','MacBook Pro',{'display':r[2],'cpu':r[4],'gpu':r[5],'storage':r[6],'ram':r[7],'color':r[8],'modelCode':r[9],'price':money(r[10])},image_files.get('MacBook Pro'),descriptions['MacBook Pro'],source_urls['Mac'])
    else:
        data=ffill_table(rr,[0,2,4,5,6,7,10])
        for r in data:
            add_product(store,'Mac','MacBook Pro',{'display':r[2],'cpu':r[4],'gpu':r[5],'ram':r[6],'storage':r[7],'color':r[8],'modelCode':r[9],'price':money(r[10]),'chip':'M5'},image_files.get('MacBook Pro'),descriptions['MacBook Pro'],source_urls['Mac'])
# Mac mini
for r in ffill_table(rows(27)[1:],[3,4,5,6,9]):
    add_product(store,'Mac','Mac mini',{'cpu':r[3],'gpu':r[4],'ram':r[5],'storage':r[6],'color':r[7],'modelCode':r[8],'price':money(r[9])},image_files.get('Mac mini'),descriptions['Mac mini'],source_urls['Mac'])
# Studio
for r in ffill_table(rows(28)[1:],[0,3,4,5,6,9]):
    add_product(store,'Mac','Mac Studio',{'cpu':r[3],'gpu':r[4],'ram':r[5],'storage':r[6],'color':r[7],'modelCode':r[8],'price':money(r[9])},image_files.get('Mac Studio'),descriptions['Mac Studio'],source_urls['Mac'])

# Accessories slide30 matrix
r30=rows(30)
model_targets=['iPhone 17','iPhone Air','iPhone 17 Pro','iPhone 17 Pro Max']
for r in ffill_table(r30[2:],[0,3]):
    fam=normalize_family(r[0])
    for idx,target in enumerate(model_targets,4):
        code=r[idx] if idx < len(r) else ''
        if code and code!='-':
            add_product(store,'Accessory',fam,{'type':'iPhone 케이스','compatible':target,'color':r[1],'modelCode':code,'price':money(r[3])},'',f'{target} 호환 케이스입니다.',source_urls['Accessory'])
# 31,32
for sn in [31,32]:
    for r in ffill_table(rows(sn)[1:],[0,1,2]):
        base=normalize_family(r[0])
        fam=f'{base} {r[1]}' if r[1] else f'{base} 케이스'
        add_product(store,'Accessory',fam,{'type':'iPhone 케이스','compatible':base,'color':r[4],'modelCode':r[5],'price':money(r[2])},'',f'{base} 호환 케이스입니다.',source_urls['Accessory'])
# 33 Pencil/keyboard
for r in ffill_table(rows(33)[1:],[0,1,4,6,7]):
    group=r[0]; subtype=r[1]; code=r[3]
    if not code: continue
    if group=='펜슬':
        fam={'프로':'Apple Pencil Pro','1세대':'Apple Pencil 1세대','USB-C':'Apple Pencil USB-C'}.get(subtype,f'Apple Pencil {subtype}')
    else:
        fam=subtype or '매직키보드'
    img=image_files.get(fam) or image_files.get('매직키보드','')
    add_product(store,'Accessory',fam,{'type':group,'color':r[5],'modelCode':code,'price':money(r[4]),'features':r[6],'compatible':r[7]},img,'호환 모델을 확인한 뒤 함께 제안할 수 있는 액세서리입니다.',source_urls['Accessory'])
# 34 folios
for r in ffill_table(rows(34,0)[2:],[0,2,3]):
    fam=normalize_family(r[0])
    for size,pc,cc in [('13인치',2,5),('11인치',3,6)]:
        if r[cc] and r[cc]!='-': add_product(store,'Accessory',fam,{'type':'iPad 케이스','size':size,'color':r[4],'modelCode':r[cc],'price':money(r[pc])},image_files.get('스마트폴리오',''),'iPad 호환 스마트 폴리오입니다.',source_urls['Accessory'])
for r in ffill_table(rows(34,1)[1:],[0,2]):
    add_product(store,'Accessory',normalize_family(r[0]),{'type':'iPad 케이스','color':r[3],'modelCode':r[4],'price':money(r[2])},image_files.get('스마트폴리오',''),'iPad mini 호환 스마트 폴리오입니다.',source_urls['Accessory'])
# watch bands
for sn in [35,36,37]:
    for r in ffill_table(rows(sn)[1:],[0,2,3,4,5]):
        if not r[6]: continue
        fam=r[3] or '워치밴드'
        add_product(store,'Accessory',fam,{'type':'Apple Watch 밴드','size':r[2],'color':r[5],'modelCode':r[6],'price':money(r[4])},image_files.get('워치밴드',''),'Apple Watch 호환 밴드입니다.',source_urls['Accessory'])
# generic accessories 38-40
for sn in [38,39,40]:
    for r in ffill_table(rows(sn)[1:],[0,4]):
        if not r[2]: continue
        group={'Mac Acc':'Mac 액세서리'}.get(r[0],r[0])
        fam=r[3]
        img=image_files.get(group,'')
        add_product(store,'Accessory',fam,{'type':group,'modelCode':r[2],'price':money(r[4])},img,'Apple 정품 액세서리입니다.',source_urls['Accessory'])

# AppleCare from iPhone section (slide 6)
for r in rows(6)[1:]:
    if not r[2] or not r[3]: continue
    covered=re.sub(r'^AppleCare\+ for ','',r[2]).replace('ProMax','Pro Max').replace('17Pro','17 Pro').replace('16Pro','16 Pro')
    fam=f'AppleCare+ {covered}'
    add_product(store,'AppleCare+',fam,{'coveredModel':covered,'modelCode':r[3],'price':money(r[4])},image_files.get('AppleCare+',''),'대상 제품과 AppleCare+ 모델코드를 함께 확인할 수 있습니다.',source_urls['AppleCare+'])

# AppleCare tables
for ti in [0,1]:
    for r in rows(42,ti)[1:]:
        if not r[0] or not r[1]: continue
        fam=f'AppleCare+ {r[0]}'
        add_product(store,'AppleCare+',fam,{'coveredModel':r[0],'modelCode':r[1],'price':money(r[2])},image_files.get('AppleCare+',''),'대상 제품과 AppleCare+ 모델코드를 함께 확인할 수 있습니다.',source_urls['AppleCare+'])

products=list(store.values())
# dedupe exact variant records and sort
for p in products:
    seen=set(); variants=[]
    for v in p['variants']:
        key=json.dumps(v,ensure_ascii=False,sort_keys=True)
        if key not in seen:
            seen.add(key); variants.append(v)
    p['variants']=variants
    prices=[v.get('price') for v in variants if isinstance(v.get('price'),int)]
    p['minPrice']=min(prices) if prices else None
    p['variantCount']=len(variants)
    if not p['description']:
        p['description']='운영 라인업의 옵션, 가격, 색상과 모델코드를 확인할 수 있습니다.'
products.sort(key=lambda p:(['iPhone','iPad','Apple Watch','AirPods','Mac','Accessory','AppleCare+'].index(p['category']),p['family']))

# AppleCare matching with category-specific rules
care=[p for p in products if p['category']=='AppleCare+']
def canon(text):
    text=str(text).lower().replace('apple watch','watch').replace('series','s').replace('(anc)','')
    return re.sub(r'[^a-z0-9가-힣]','',text)
def generation_priority(text):
    t=str(text).upper()
    for gen,score in [('M5',50),('M4',40),('M3',30),('M2',20),('S11',15),('SE 3',14),('ULTRA 3',13)]:
        if gen in t: return score
    return 0
for p in products:
    if p['category']=='AppleCare+': continue
    pf=canon(p['family'])
    matches=[]
    for c in care:
        v=c['variants'][0]
        covered=v.get('coveredModel','')
        cm=canon(covered)
        ok=False
        if p['category']=='iPhone':
            ok=(cm==pf)
        elif p['category']=='AirPods':
            target=canon(p['family'].replace('(ANC)',''))
            ok=(cm==target)
        elif p['category']=='Apple Watch':
            target={'Apple Watch SE3':'Watch SE 3','Apple Watch Series 11':'Watch S11','Apple Watch Ultra3':'Watch Ultra 3'}.get(p['family'],p['family'])
            ok=(cm==canon(target))
        elif p['category']=='iPad':
            ok=cm.startswith(pf)
        elif p['category']=='Mac':
            if p['family']=='iMac': ok=cm.startswith('imac')
            elif p['family']=='Mac mini': ok=cm.startswith('macmini')
            elif p['family']=='MacBook Neo': ok=cm.startswith('macbookneo')
            elif p['family']=='MacBook Air': ok=cm.startswith('macbookair')
            elif p['family']=='MacBook Pro': ok=cm.startswith('macbookpro')
        if ok:
            matches.append({'family':c['family'],'variant':v,'priority':generation_priority(covered)})
    if matches:
        matches.sort(key=lambda x:(-x['priority'],x['family']))
        for x in matches: x.pop('priority',None)
        p['appleCareOptions']=matches

# Data validation flags from the source PPT
code_locations=defaultdict(list)
for p in products:
    for v in p['variants']:
        code=v.get('modelCode','')
        if code: code_locations[code].append(p['family'])
validation_warnings=[]
for p in products:
    for v in p['variants']:
        code=v.get('modelCode','')
        warnings=[]
        if code and len(code_locations[code])>1:
            warnings.append('동일 모델코드가 운영표의 다른 항목에도 사용됨')
        if code and not re.match(r'^[A-Z0-9]+/[A-Z]+$',code):
            warnings.append('모델코드 표기 형식 확인 필요')
        if warnings:
            v['warning']=' / '.join(warnings)
            validation_warnings.append({'family':p['family'],'modelCode':code,'warning':v['warning']})

meta={
 'generatedAt':'2026-07-16',
 'source':'2607)Apple_운영_라인업 (1).pptx',
 'productCount':len(products),
 'variantCount':sum(len(p['variants']) for p in products),
 'categoryCounts':dict(defaultdict(int))
}
for p in products: meta['categoryCounts'][p['category']]=meta['categoryCounts'].get(p['category'],0)+1
meta['validationWarningCount']=len(validation_warnings)
meta['validationWarnings']=validation_warnings
(ROOT/'data.json').write_text(json.dumps({'meta':meta,'products':products},ensure_ascii=False,indent=2),encoding='utf-8')
print(json.dumps(meta,ensure_ascii=False,indent=2))
print('assets',len(list(ASSET.glob('*'))))
